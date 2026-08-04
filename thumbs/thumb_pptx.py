#!/usr/bin/env python3
"""
サムネイルを編集可能な pptx（LibreOffice Impress用）として書き出す。

    python3 thumbs/thumb_pptx.py 3

thumbs/pptx/epNNN.pptx が生成される。バッジ・タイトル・脚注・選択肢ピルは
動かさない背景画像として焼き込み、立ち絵とセリフ（吹き出し）だけを
ドラッグで動かせる別要素として配置する。

Impress で開いて位置を微調整したら、
    ファイル → エクスポート → 画像としてエクスポート
    幅 1280px 高さ 670px を指定して PNG 書き出し
    → thumbs/out/epNNN.png に上書き保存

または、書き出したファイルを引数に thumbs/pptx_to_png.py で自動変換できる。
"""
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from thumb import ROOT, H, W, build, read_meta  # noqa: E402

EMU_PER_PX = 9525  # 96dpi換算（PILのピクセルとpptxの実寸を一致させる）
FONT_BLACK = 'Noto Sans CJK JP Black'
FONT_MED = 'Noto Sans CJK JP Medium'
FONT_REG = 'Noto Sans CJK JP'


def px(v):
    return Emu(int(round(v * EMU_PER_PX)))


def pt_from_px(v):
    # PILのフォントサイズ(px相当)を96dpi基準でpt換算
    return Pt(round(v / 1.3333, 1))


def rgb(t):
    return RGBColor(*t)


def add_pill(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px(x), px(y), px(w), px(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_rgb)
    if line_rgb:
        shp.line.color.rgb = rgb(line_rgb)
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5  # カプセル型に丸める
    except IndexError:
        pass
    return shp


def set_text(shp, text, font, size_px, color, bold, align, word_wrap=True):
    tf = shp.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = px(18)
    tf.margin_right = px(18)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = pt_from_px(size_px)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = rgb(color)


def main():
    if len(sys.argv) < 2:
        sys.exit('使い方: python3 thumbs/thumb_pptx.py <話数>')
    ep = int(sys.argv[1])
    meta = read_meta(ep)

    bg_img, layout = build(ep, meta, bake_char=False, bake_chat=False)
    out_dir = ROOT / 'thumbs' / 'pptx'
    out_dir.mkdir(parents=True, exist_ok=True)
    bg_path = out_dir / f'ep{ep:03d}_bg.png'
    bg_img.save(bg_path, 'PNG')

    prs = Presentation()
    prs.slide_width = px(W)
    prs.slide_height = px(H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 白紙レイアウト

    # 動かさない背景（バッジ・タイトル・脚注・選択肢ピル）
    slide.shapes.add_picture(str(bg_path), 0, 0, width=px(W), height=px(H))

    # 立ち絵（ドラッグで移動・ハンドルでリサイズ可）
    for meta_key, box_key in (('thumb_char', 'char_box'), ('thumb_char2', 'char_box2')):
        if layout[box_key] and meta.get(meta_key):
            char_path = ROOT / meta[meta_key]
            cx, cy, cw, ch = layout[box_key]
            if char_path.exists():
                slide.shapes.add_picture(str(char_path), px(cx), px(cy), width=px(cw), height=px(ch))
            else:
                print(f'! 立ち絵が見つかりません: {char_path}')

    # セリフ（名前チップ＋吹き出し。ドラッグで移動可）
    for row in layout['rows']:
        cx, cy, cw, ch = row['chip_box']
        chip = add_pill(slide, cx, cy, cw, ch, row['name_col'])
        set_text(chip, row['name'], FONT_BLACK, 24, (255, 255, 255), True, PP_ALIGN.CENTER,
                 word_wrap=False)

        bx, by, bw, bh = row['bubble_box']
        bubble = add_pill(slide, bx, by, bw, bh, row['bg'], row['ln_col'])
        set_text(bubble, row['text'], FONT_MED, 27, (31, 35, 55), False, PP_ALIGN.LEFT)

    if layout['sub']:
        sx, sy = layout['sub']['pos']
        box = slide.shapes.add_textbox(px(sx), px(sy), px(W - sx - 20), px(34))
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = layout['sub']['text']
        run.font.size = pt_from_px(19)
        run.font.name = FONT_REG
        run.font.color.rgb = rgb((95, 103, 130))

    out_path = out_dir / f'ep{ep:03d}.pptx'
    prs.save(out_path)
    print(f'生成しました: {out_path}')
    print('LibreOffice Impress で開いて、立ち絵とセリフの位置を調整してください。')
    print('  soffice thumbs/pptx/ep%03d.pptx' % ep)
    print('調整後は次のコマンドで PNG に書き戻せます:')
    print('  python3 thumbs/pptx_to_png.py %d' % ep)


if __name__ == '__main__':
    main()
